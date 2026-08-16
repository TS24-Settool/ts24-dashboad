#!/usr/bin/env python3
"""TS24 Suspension Report v2 generator.

Workbench の `PhaseRunCompareWidget` から呼ばれ、`lap_suspension`（正本DB read-only）を
ソースに Braking / Apex / Exit の姿勢・ダンピング速度を PowerPoint (.pptx) で出力する。

Report v2（2026-07-02 / Tatsuki feedback 反映・設計 = reports/workbench_report_v2_design_20260702.md）:
  1. グラフのラベル/凡例をプロット外へ（matplotlib Agg・被り回避）
  2. Lap time は M:SS,CC（例 103.739 -> "1:43,74"）
  3. 表はヘッダ2行＋単位＋Braking/Apex/Exit セル色分け＋説明行
  4. 0% を Missing / Coverage / Structural n/a に分離明示
  5. Run内 Lap by lap 分析ページ（time / position / speed progression + run detail）
  6. Braking=red / Apex=blue / Exit=green を全所で統一

DB は `file:...?mode=ro` の read-only 接続のみ。schema 変更・書込は一切しない。
matplotlib / python-pptx が未導入の場合は ReportUnavailableError を送出し、呼び出し側で
アプリを落とさず message box 表示できるようにする。
"""
from __future__ import annotations

import math
import sqlite3
from datetime import datetime
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_DB = SCRIPT_DIR.parent / "02_DATABASE" / "ts24_unified.db"
DEFAULT_OUT = SCRIPT_DIR / "reports" / "pptx"

LAP_MIN, LAP_MAX = 60.0, 300.0          # 有効ラップ（アウトラップ/計測エラー除外）
F_MAX, R_MAX = 130.0, 70.0              # Full Stroke（物理限界）
AVG_NMIN = 5                            # avg 信頼下限（§44）
RUN_CHART_CAP = 6                       # 1グラフに重ねる Run 上限（凡例肥大回避）
RUN_DETAIL_CAP = 6                      # Run detail ページ上限

# フェーズ色（Workbench _PHASE_COLORS と一致）
PHASE_COLORS = {"Braking": "#C0392B", "Apex": "#0078D4", "Exit": "#2E9E4F"}
PHASE_FILL = {"Braking": "FBE9E7", "Apex": "E8F1FB", "Exit": "E9F6EE"}   # 表セル薄色
PHASES = ["Braking", "Apex", "Exit"]

# フェーズ×side の Position 列（DataFrame は全小文字）
PHASE_POS = {
    "Braking": ("brk_susf_avg", "brk_susr_avg"),
    "Apex":    ("apex_susf_avg", "apex_susr_avg"),
    "Exit":    ("ce_susf_avg",   "ce_susr_avg"),
}
# フェーズ×side の Suspension Speed 列（avg, peak, 短縮タグ）— §44 本命方向
PHASE_SPD = {
    "Braking": {"F": ("brk_f_dive_spd_avg", "brk_f_dive_spd_peak", "F-Dive"),
                "R": ("brk_r_reb_spd_avg",  "brk_r_reb_spd_peak",  "R-Reb")},
    "Apex":    {"F": ("apex_f_dive_spd_avg", "apex_f_dive_spd_peak", "F-Dive"),
                "R": ("apex_r_dive_spd_avg", "apex_r_dive_spd_peak", "R-Dive")},
    "Exit":    {"F": ("ce_f_reb_spd_avg", "ce_f_reb_spd_peak", "F-Reb"),
                "R": ("ce_r_spd_avg",     "ce_r_spd_peak",     "R|v|")},
}
RUN_PALETTE = ["#0078D4", "#FF8C00", "#2E9E4F", "#C0392B", "#8E44AD",
               "#16A085", "#D4A017", "#E91E63", "#3F51B5", "#5D6D7E"]

SPEED_NOTE = "Susp speed = relative damping-speed index (mm/s, uncalibrated) — NOT vehicle speed (km/h)"

SLOW_LAP_FACTOR = 1.07   # session median × factor 超のラップ = slow outlier（report-only filter・§66）

# ── Report Update 2026-07-10（presentation/readability・report-only）──
OUTLIER_IQR_K = 1.5      # IQR robust 外れ値係数（report-only visual flag・DB/抽出 無変更）
OUTLIER_LABEL_CAP = 6    # 外れ値ラベル上限 / パネル（クラッタ回避 cap）
TREND_OUTLIER_NOTE = (
    "All valid laps of ALL selected runs after the page-2 lap filter - no extra filtering. "
    "Metric = Front phase position (mm), single clearest family; rear = phase summary pages. "
    "Red ring = report-only outlier flag (IQR rule: outside Q1/Q3 +/-1.5xIQR per phase, "
    f"max {OUTLIER_LABEL_CAP} labels); laps are NOT removed - no DB/extraction change.")
DIST_NOTE = (
    "Box = Q1-Q3 + median per run; dots = individual laps (after page-2 lap filter). "
    "Red ring = report-only lap-time outlier flag (IQR rule across all shown laps, "
    f"max {OUTLIER_LABEL_CAP} labels); star = fastest lap. No DB/extraction change.")

# フェーズ毎の速度サンプリング窓（§66 audit）。avg はこの窓内の MEAN velocity。
PHASE_SPEED_REGION = {"Braking": "deep-stroke / settled",
                      "Apex":    "mid-stroke",
                      "Exit":    "corner-exit (sparse)"}
# §66: brk_f_dive_spd_avg は FULL_BRAKING(SUSP_FRONT 90-130mm=フォークが既に深く沈んだ準定常)内の
# 平均速度で、制動初期の高速ダイブ過渡(SUSP_FRONT<90)は含まない → Apex>Braking を「apex で速く沈む」と
# 読ませないための注記。相対 index であり車速 km/h とは別物。
SPEED_WINDOW_NOTE = ("F/R speed = MEAN suspension velocity WITHIN each phase window "
                     "(Braking = deep/settled fork 90-130mm; Apex = mid-stroke 50-100mm), "
                     "NOT the peak brake dive-in rate. Braking avg is low because the fork is already "
                     "near bottom — do NOT read Apex>Braking as 'the front dives faster at apex'.")
PEAK_NOTE = ("avg needs n>=5 laps. Peak columns are not shown/compared across phases: "
             "brk_f_dive_spd_peak is a legacy MAX while other *_peak are p95 (not directly comparable).")
OHLINS_NOTE = ("This index is NOT directly comparable to Ohlins low/high-speed compression/rebound "
               "(force-vs-shaft-velocity) settings — it is an observed travel-rate, not a damper transfer function.")


class ReportUnavailableError(RuntimeError):
    """matplotlib / python-pptx 未導入など、Report 生成不能時に送出。"""


# ── フォーマッタ ───────────────────────────────────────────────────────
def format_lap_time(sec) -> str:
    """秒 -> "M:SS,CC"（欧州式カンマ小数・センチ秒）。例 103.739 -> "1:43,74"。

    None / NaN / 非正 は "n/a"。59.995 の繰り上げをガード。
    """
    try:
        s = float(sec)
    except (TypeError, ValueError):
        return "n/a"
    if s != s or s <= 0:            # NaN or <=0
        return "n/a"
    m = int(s) // 60
    cs = round((s - m * 60) * 100)  # センチ秒
    if cs >= 6000:                  # 59.995 -> 次の分へ
        m += 1
        cs -= 6000
    return f"{m}:{cs // 100:02d},{cs % 100:02d}"


def _num(v):
    try:
        f = float(v)
    except (TypeError, ValueError):
        return None
    return None if f != f else f


def _mean(vals):
    xs = [x for x in (_num(v) for v in vals) if x is not None]
    return sum(xs) / len(xs) if xs else None


def _median(vals):
    xs = sorted(x for x in (_num(v) for v in vals) if x is not None)
    if not xs:
        return None
    n = len(xs)
    mid = n // 2
    return xs[mid] if n % 2 else (xs[mid - 1] + xs[mid]) / 2


def _iqr_bounds(vals, k=OUTLIER_IQR_K):
    """IQR robust bounds（report-only 外れ値 flag 用・指示書 2026-07-10）。

    有効値 < 4 は None（少数サンプルで flag しない）。データからの除外はしない。
    """
    xs = sorted(x for x in (_num(v) for v in vals) if x is not None)
    if len(xs) < 4:
        return None

    def _pct(p):
        i = (len(xs) - 1) * p
        lo, hi = int(math.floor(i)), int(math.ceil(i))
        return xs[lo] + (xs[hi] - xs[lo]) * (i - lo)

    q1, q3 = _pct(0.25), _pct(0.75)
    iqr = q3 - q1
    return q1 - k * iqr, q3 + k * iqr


def run_short_label(rec) -> str:
    """Run 短縮ラベル: "DA77 SP R1"。"""
    rider = rec.get("rider", "?")
    sess = rec.get("session", "?")
    try:
        rn = f"R{int(rec.get('run_no'))}"
    except (TypeError, ValueError):
        rn = f"R{rec.get('run_no', '?')}"
    return f"{rider} {sess} {rn}"


# ── データ取得（read-only）────────────────────────────────────────────
def load_lap_suspension(db_path=DEFAULT_DB, circuit=None, rider=None,
                        session=None, run_ids=None):
    """`lap_suspension` を read-only(mode=ro) で読み、フィルタ済み DataFrame を返す。

    `laps.is_outlap` を lap_id JOIN で付与（out/in ラップ除外の任意強化に使用）。
    """
    import pandas as pd

    uri = f"file:{Path(db_path)}?mode=ro"
    with sqlite3.connect(uri, uri=True) as con:
        con.row_factory = sqlite3.Row
        rows = con.execute(
            "SELECT ls.*, l.is_outlap AS _is_outlap "
            "FROM lap_suspension ls "
            "LEFT JOIN laps l ON l.lap_id = ls.lap_id"
        ).fetchall()
    df = pd.DataFrame([dict(r) for r in rows])
    if df.empty:
        return df
    df.columns = [c.lower() for c in df.columns]
    if circuit and "circuit" in df.columns:
        df = df[df["circuit"] == circuit]
    if rider and "rider" in df.columns:
        df = df[df["rider"] == rider]
    if session and "session" in df.columns:
        df = df[df["session"] == session]
    if "lap_time_s" in df.columns:
        df = df[df["lap_time_s"].between(LAP_MIN, LAP_MAX)]
    if run_ids:
        df = df[df["run_id"].isin(list(run_ids))]
    return df.reset_index(drop=True)


def _run_records(df):
    """Run 単位のメタ（run_id 昇順の run_no でソート）。"""
    recs = []
    for run_id, g in df.groupby("run_id"):
        r0 = g.iloc[0]
        recs.append({
            "run_id": run_id,
            "rider": r0.get("rider"), "session": r0.get("session"),
            "circuit": r0.get("circuit"), "round": r0.get("round"),
            "run_no": r0.get("run_no"), "n_laps": len(g),
            "_df": g,
        })
    recs.sort(key=lambda r: (str(r.get("session")), _num(r.get("run_no")) or 0))
    return recs


# ── 集計 ───────────────────────────────────────────────────────────────
def session_summary(df) -> dict:
    laps = df["lap_time_s"].tolist() if "lap_time_s" in df.columns else []
    return {
        "n_runs": df["run_id"].nunique() if "run_id" in df.columns else 0,
        "n_laps": len(df),
        "best": min([x for x in (_num(v) for v in laps) if x is not None], default=None),
        "median": _median(laps),
        "mean": _mean(laps),
    }


def run_best(g) -> float | None:
    """Run 内 best lap（有効ラップ最小・is_outlap!=1 を優先）。"""
    import pandas as pd
    valid = g
    if "_is_outlap" in g.columns:
        v2 = g[g["_is_outlap"] != 1]
        if not v2.empty:
            valid = v2
    xs = [x for x in (_num(v) for v in valid["lap_time_s"]) if x is not None]
    return min(xs) if xs else None


def phase_run_stats(df, phase) -> list:
    """Run 別に position(F/R avg) と speed(F/R avg) を集計。"""
    fcol, rcol = PHASE_POS[phase]
    spd = PHASE_SPD[phase]
    out = []
    for rec in _run_records(df):
        g = rec["_df"]
        out.append({
            "label": run_short_label(rec), "run_id": rec["run_id"],
            "f_pos": _mean(g[fcol]) if fcol in g else None,
            "r_pos": _mean(g[rcol]) if rcol in g else None,
            "f_spd": _mean(g[spd["F"][0]]) if spd["F"][0] in g else None,
            "r_spd": _mean(g[spd["R"][0]]) if spd["R"][0] in g else None,
        })
    return out


def lap_series(g):
    """Run 内 lap 昇順の系列（lap_no, lap_time, best差, phase pos/spd）。"""
    g = g.sort_values("lap_no")
    best = run_best(g)
    laps, times, delta = [], [], []
    for _, row in g.iterrows():
        ln = _num(row.get("lap_no"))
        lt = _num(row.get("lap_time_s"))
        if ln is None:
            continue
        laps.append(int(ln))
        times.append(lt)
        delta.append((lt - best) if (lt is not None and best is not None) else None)
    return {"lap_no": laps, "lap_time": times, "delta": delta, "best": best, "_df": g}


def data_quality(df) -> list:
    """指標グループ別に populated / total / missing% / 意味を返す。"""
    total = len(df)
    groups = [
        ("Braking position (F/R)", ["brk_susf_avg", "brk_susr_avg"], "sensor position at full braking"),
        ("Apex position (F/R)",    ["apex_susf_avg", "apex_susr_avg"], "sensor position at apex"),
        ("Exit position (F/R)",    ["ce_susf_avg", "ce_susr_avg"], "sensor position at corner exit"),
        ("Braking speed (F dive/R reb)", ["brk_f_dive_spd_avg", "brk_r_reb_spd_avg"], "damping speed index"),
        ("Apex speed (F/R)",       ["apex_f_dive_spd_avg", "apex_r_dive_spd_avg"], "damping speed index"),
        ("Exit speed (F reb/R)",   ["ce_f_reb_spd_avg", "ce_r_spd_avg"], "damping speed index (Exit sparse)"),
    ]
    rows = []
    for name, cols, meaning in groups:
        present = [c for c in cols if c in df.columns]
        if not present:
            rows.append({"group": name, "populated": 0, "total": total,
                         "missing_pct": None, "note": "n/a (column not available)"})
            continue
        # ラップ単位で「いずれかの列が非NULL」を populated とする
        pop = 0
        for _, row in df[present].iterrows():
            if any(_num(row[c]) is not None for c in present):
                pop += 1
        miss = (total - pop) / total * 100 if total else None
        structural = "Exit" in name
        if miss == 0:
            note = f"Missing 0% (all {total} laps populated) — {meaning}"
        elif structural:
            note = f"Missing {miss:.0f}% — structural: sparse CORNER_EXIT — {meaning}"
        else:
            note = f"Missing {miss:.0f}% ({pop}/{total} populated) — {meaning}"
        rows.append({"group": name, "populated": pop, "total": total,
                     "missing_pct": miss, "note": note})
    return rows


def apply_lap_filter(df, enabled=True):
    """Report-only の決定論的 slow/out-lap フィルタ（§66 Tier1・DB書込なし）。

    除外規則: (1) out/in ラップ（`_is_outlap==1`・列があるときのみ）
              (2) session ごとの有効ラップ中央値 × SLOW_LAP_FACTOR を超える slow lap。
    返り値 (kept_df, excluded) — excluded=[{label, lap_no, lap_time, reason}]。
    全ラップが除外される退化時はフィルタ無効化（元 df をそのまま返す）。
    表示（適用 filter・除外 lap 一覧・理由）は必ず Data Quality ページで開示する。
    """
    if not enabled or df is None or df.empty or "lap_time_s" not in df.columns:
        return df, []
    has_outlap = "_is_outlap" in df.columns
    grp_col = "session" if "session" in df.columns else "run_id"
    sess_med = {}
    for key, g in df.groupby(grp_col):
        base = g[g["_is_outlap"] != 1] if has_outlap else g
        sess_med[key] = _median(base["lap_time_s"].tolist())
    keep_idx, excluded = [], []
    for idx, row in df.iterrows():
        lt = _num(row.get("lap_time_s"))
        is_out = has_outlap and _num(row.get("_is_outlap")) == 1
        med = sess_med.get(row.get(grp_col))
        is_slow = lt is not None and med is not None and lt > med * SLOW_LAP_FACTOR
        if is_out or is_slow:
            excluded.append({
                "label": run_short_label({"rider": row.get("rider"),
                                          "session": row.get("session"),
                                          "run_no": row.get("run_no")}),
                "lap_no": _num(row.get("lap_no")), "lap_time": lt,
                "reason": "out/in lap" if is_out else f"slow (> median x{SLOW_LAP_FACTOR:.2f})"})
        else:
            keep_idx.append(idx)
    if not keep_idx:                        # 退化ガード: 全除外はしない
        return df, []
    return df.loc[keep_idx].reset_index(drop=True), excluded


def lap_filter_note(excluded, enabled=True, cap=6):
    """Data Quality ページ用の lap-filter 開示テキスト（英語・page-2 開示）。"""
    if not enabled:
        return "Lap filter: OFF — all laps shown (unfiltered)."
    if not excluded:
        return (f"Lap filter (report-only, no DB change): ON — 0 laps excluded "
                f"(no out-laps; none slower than session median x{SLOW_LAP_FACTOR:.2f}).")
    items = []
    for e in excluded[:cap]:
        ln = f"L{int(e['lap_no'])}" if e.get("lap_no") is not None else "L?"
        items.append(f"{e['label']} {ln} {format_lap_time(e['lap_time'])} ({e['reason']})")
    more = f"  [+{len(excluded) - cap} more]" if len(excluded) > cap else ""
    return ("Lap filter (report-only, no DB change): ON — "
            f"excluded {len(excluded)} lap(s): " + "; ".join(items) + more)


# ── matplotlib チャート（Agg・ラベルはプロット外）─────────────────────
def _mpl():
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.ticker import FuncFormatter
        return plt, FuncFormatter
    except Exception as exc:  # pragma: no cover
        raise ReportUnavailableError(f"matplotlib unavailable: {exc}") from exc


def _style(plt):
    plt.rcParams.update({
        "figure.dpi": 150, "savefig.dpi": 150,
        "font.size": 11, "axes.titlesize": 13, "axes.titleweight": "bold",
        "axes.grid": True, "grid.alpha": 0.3, "grid.linestyle": "--",
        "axes.spines.top": False, "axes.spines.right": False,
        "figure.autolayout": False,
    })


def _save(fig, tmpdir, name, tight=True):
    p = Path(tmpdir) / f"{name}.png"
    fig.savefig(p, bbox_inches=("tight" if tight else None), facecolor="white")
    import matplotlib.pyplot as plt
    plt.close(fig)
    return p


def _bar_value_labels(ax, bars, vals, fmt, fontsize=8):
    """棒の上に実数値ラベル（指示書 2026-07-10 §1）。

    欠損(None/NaN)はラベルを付けない（0 と表示しない）。棒が多い場合は縮小
    フォントでクラッタ回避し、上方向に headroom を追加して軸/タイトルと非重複。
    """
    if len(vals) > 12:
        fontsize = 6.5
    any_label = False
    for rect, v in zip(bars, vals):
        f = _num(v)
        if f is None:
            continue
        ax.annotate(fmt(f), (rect.get_x() + rect.get_width() / 2, rect.get_height()),
                    ha="center", va="bottom", fontsize=fontsize,
                    xytext=(0, 2), textcoords="offset points")
        any_label = True
    if any_label:
        lo, hi = ax.get_ylim()
        ax.set_ylim(lo, hi + (hi - lo) * 0.10)


def chart_run_overview(df, tmpdir):
    """Run 別 best/median lap 棒グラフ（y軸 = M:SS,CC・値ラベルは棒の外）。"""
    plt, FuncFormatter = _mpl()
    recs = _run_records(df)
    labels = [run_short_label(r) for r in recs]
    best = [run_best(r["_df"]) for r in recs]
    med = [_median(r["_df"]["lap_time_s"]) for r in recs]
    x = range(len(labels))
    fig, ax = plt.subplots(figsize=(11, 5.2))
    w = 0.38
    b1 = ax.bar([i - w / 2 for i in x], [b or 0 for b in best], w,
                label="Best lap", color="#0078D4")
    b2 = ax.bar([i + w / 2 for i in x], [m or 0 for m in med], w,
                label="Median lap", color="#9DC3E6")
    lo = min([v for v in best + med if v], default=0)
    hi = max([v for v in best + med if v], default=1)
    ax.set_ylim(max(0, lo - (hi - lo) * 0.25), hi + (hi - lo) * 0.15)
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: format_lap_time(v)))
    ax.set_ylabel("Lap time (M:SS,CC)")
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
    for bars, vals in ((b1, best), (b2, med)):
        for rect, v in zip(bars, vals):
            if v:
                ax.annotate(format_lap_time(v), (rect.get_x() + rect.get_width() / 2,
                            rect.get_height()), ha="center", va="bottom",
                            fontsize=8, xytext=(0, 2), textcoords="offset points")
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.28), ncol=2, frameon=False)
    ax.set_title("Run Overview — Best / Median Lap")
    return _save(fig, tmpdir, "run_overview")


def chart_phase_summary(df, phase, tmpdir):
    """1フェーズの Run 別: F position / R position（独立Y=small multiples）と F/R damping speed。

    §66 feedback①: F と R の position を別パネル・独立Y にして、F の大レンジで R が潰れないようにする。
    """
    plt, _ = _mpl()
    stats = phase_run_stats(df, phase)
    labels = [s["label"] for s in stats]
    x = range(len(labels))
    color = PHASE_COLORS[phase]
    fig, (axpf, axpr, axs) = plt.subplots(1, 3, figsize=(13.2, 5.0))
    # F/R position を別パネル（独立軸）に分離 → R が F のスケールに潰れない
    bpf = axpf.bar(list(x), [s["f_pos"] or 0 for s in stats], 0.6, color=color)
    axpf.set_ylabel("Position (mm)")
    axpf.set_title(f"{phase} — F position (mm)")
    _bar_value_labels(axpf, bpf, [s["f_pos"] for s in stats], lambda v: f"{v:.1f}")
    bpr = axpr.bar(list(x), [s["r_pos"] or 0 for s in stats], 0.6, color=color, alpha=0.55)
    axpr.set_ylabel("Position (mm)")
    axpr.set_title(f"{phase} — R position (mm)")
    _bar_value_labels(axpr, bpr, [s["r_pos"] for s in stats], lambda v: f"{v:.1f}")
    ftag, rtag = PHASE_SPD[phase]["F"][2], PHASE_SPD[phase]["R"][2]
    ws = 0.38
    bsf = axs.bar([i - ws / 2 for i in x], [s["f_spd"] or 0 for s in stats], ws,
                  label=f"F {ftag}", color=color)
    bsr = axs.bar([i + ws / 2 for i in x], [s["r_spd"] or 0 for s in stats], ws,
                  label=f"R {rtag}", color=color, alpha=0.45)
    _bar_value_labels(axs, bsf, [s["f_spd"] for s in stats], lambda v: f"{v:.0f}", fontsize=7)
    _bar_value_labels(axs, bsr, [s["r_spd"] for s in stats], lambda v: f"{v:.0f}", fontsize=7)
    axs.set_ylabel("Susp speed (idx, mm/s uncal.)")
    axs.set_title(f"{phase} — Damping speed avg\n({PHASE_SPEED_REGION[phase]})")
    axs.legend(loc="upper center", bbox_to_anchor=(0.5, -0.30), ncol=2, frameon=False)
    for ax in (axpf, axpr, axs):
        ax.set_xticks(list(x))
        ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=8)
    fig.suptitle(f"{phase} Phase Summary", fontsize=14, fontweight="bold", color=color)
    fig.subplots_adjust(wspace=0.42, bottom=0.26, top=0.84)
    return _save(fig, tmpdir, f"phase_{phase.lower()}")


def _capped_runs(df):
    recs = _run_records(df)
    capped = recs[:RUN_CHART_CAP]
    return capped, max(0, len(recs) - RUN_CHART_CAP)


def chart_lap_time_progression(df, tmpdir):
    """Lap by lap: lap time 推移（Run 毎の線・best マーカー・y=M:SS,CC）。"""
    plt, FuncFormatter = _mpl()
    recs, extra = _capped_runs(df)
    fig, ax = plt.subplots(figsize=(11.5, 5.4))
    for i, rec in enumerate(recs):
        s = lap_series(rec["_df"])
        col = RUN_PALETTE[i % len(RUN_PALETTE)]
        xs = [ln for ln, t in zip(s["lap_no"], s["lap_time"]) if t is not None]
        ys = [t for t in s["lap_time"] if t is not None]
        if not xs:
            continue
        ax.plot(xs, ys, "-o", color=col, ms=4, lw=1.6, label=run_short_label(rec))
        if s["best"] is not None:
            bx = xs[ys.index(min(ys))]
            ax.plot([bx], [min(ys)], marker="*", color=col, ms=14,
                    markeredgecolor="black", markeredgewidth=0.5, zorder=5)
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: format_lap_time(v)))
    ax.set_xlabel("Lap No")
    ax.set_ylabel("Lap time (M:SS,CC)")
    ax.set_title("Lap-by-lap — Lap Time Progression  (★ = run best)")
    ax.legend(loc="center left", bbox_to_anchor=(1.01, 0.5), frameon=False, fontsize=9)
    note = f"+{extra} more run(s) omitted — see table" if extra else ""
    if note:
        ax.annotate(note, (1.0, -0.12), xycoords="axes fraction", ha="right",
                    fontsize=8, color="#888")
    fig.subplots_adjust(right=0.78)
    return _save(fig, tmpdir, "lap_time_prog")


def chart_all_laps_phase_trend(df, tmpdir):
    """All Laps Phase Trend & Outliers（指示書 2026-07-10 §2）。

    3フェーズパネル（Braking/Apex/Exit）。X = 連番ラップシーケンス（run 昇順・
    run 内 lap_no 昇順）。色 = run・マーカー = 各 lap・破線 = run 毎の median。
    選択された ALL run / lap filter 後の ALL valid lap を表示（新規 silent filter 禁止・cap なし）。
    metric = F phase position（単一の最明瞭 family・ページ注記で明示）。
    外れ値 = IQR robust rule（per phase）の report-only visual flag（赤リング + 短ラベル・
    ラベルは OUTLIER_LABEL_CAP まで）。データからは除外しない。
    """
    plt, _ = _mpl()
    recs = _run_records(df)              # 全 run — RUN_CHART_CAP は適用しない
    run_pts = []
    xseq = 1
    for rec in recs:
        g = rec["_df"].sort_values("lap_no")
        pts = []
        for _, row in g.iterrows():
            ln = _num(row.get("lap_no"))
            if ln is None:
                continue
            pts.append((xseq, int(ln), row))
            xseq += 1
        run_pts.append((rec, pts))
    fig, axes = plt.subplots(1, 3, figsize=(13.2, 5.6), sharex=True)
    for pi, phase in enumerate(PHASES):
        ax = axes[pi]
        color = PHASE_COLORS[phase]
        fcol = PHASE_POS[phase][0]
        allv = []
        for _rec, pts in run_pts:
            allv.extend(v for v in (_num(r.get(fcol)) for _x, _ln, r in pts)
                        if v is not None)
        bounds = _iqr_bounds(allv)
        outliers = []                    # (distance, x, value, label)
        has = False
        for i, (rec, pts) in enumerate(run_pts):
            c = RUN_PALETTE[i % len(RUN_PALETTE)]
            xs, ys, lns = [], [], []
            for _x, _ln, row in pts:
                v = _num(row.get(fcol))
                if v is None:
                    continue
                xs.append(_x)
                ys.append(v)
                lns.append(_ln)
            if not xs:
                continue
            has = True
            ax.plot(xs, ys, "-o", color=c, lw=1.0, ms=4, alpha=0.8)
            med = _median(ys)
            if med is not None:
                ax.plot([xs[0], xs[-1]], [med, med], ls="--", lw=1.1,
                        color=c, alpha=0.45)
            if bounds:
                lo, hi = bounds
                try:
                    rtag = f"R{int(rec.get('run_no'))}"
                except (TypeError, ValueError):
                    rtag = f"R{rec.get('run_no')}"
                for _x, v, ln in zip(xs, ys, lns):
                    if v < lo or v > hi:
                        d = (lo - v) if v < lo else (v - hi)
                        outliers.append((d, _x, v, f"{rtag} L{ln} {v:.1f}"))
        if not has:
            ax.text(0.5, 0.5, "n/a", ha="center", va="center",
                    transform=ax.transAxes, color="#aaa", fontsize=12)
        if outliers:
            ax.scatter([o[1] for o in outliers], [o[2] for o in outliers], s=140,
                       facecolors="none", edgecolors="#D62728", linewidths=1.8,
                       zorder=6)
            outliers.sort(key=lambda o: -o[0])   # 極端な順にラベル（cap あり）
            for k, (_d, ox, ov, txt) in enumerate(outliers[:OUTLIER_LABEL_CAP]):
                dy = 9 if k % 2 == 0 else -14
                ax.annotate(txt, (ox, ov), xytext=(0, dy),
                            textcoords="offset points", ha="center", fontsize=7,
                            color="#D62728", fontweight="bold", zorder=7)
            if len(outliers) > OUTLIER_LABEL_CAP:
                ax.annotate(f"+{len(outliers) - OUTLIER_LABEL_CAP} more flagged",
                            (0.99, 0.02), xycoords="axes fraction", ha="right",
                            fontsize=7, color="#D62728")
        lo_y, hi_y = ax.get_ylim()
        ax.set_ylim(lo_y, hi_y + (hi_y - lo_y) * 0.10)   # ラベル headroom
        ax.set_title(f"{phase} — F position (mm)", fontsize=11, color=color)
        ax.set_xlabel("Lap sequence")
    axes[0].set_ylabel("F position (mm)")
    handles = [plt.Line2D([0], [0], color=RUN_PALETTE[i % len(RUN_PALETTE)], lw=2,
               marker="o", ms=4, label=run_short_label(rec))
               for i, (rec, _p) in enumerate(run_pts)]
    fig.legend(handles=handles, loc="lower center", ncol=min(max(len(handles), 1), 6),
               frameon=False, fontsize=8, bbox_to_anchor=(0.5, -0.02))
    fig.suptitle("All Laps Phase Trend & Outliers", fontsize=14, fontweight="bold",
                 y=0.985)
    # PDF 経路にも規則開示が載るよう、注記を図内にも焼き込む（PPTX は slide note 併用）
    fig.text(0.5, 0.925,
             "All valid laps of all selected runs (after page-2 lap filter, no extra filtering) "
             "- Front position family. Red ring = report-only IQR outlier flag "
             "(Q1/Q3 +/-1.5xIQR per phase); no DB/extraction change.",
             ha="center", fontsize=8, style="italic", color="#666666")
    fig.subplots_adjust(wspace=0.24, bottom=0.26, top=0.84)
    return _save(fig, tmpdir, "all_laps_phase_trend")


def chart_lap_time_distribution(df, tmpdir):
    """Lap Time Distribution（指示書 2026-07-10 §3）。

    run 別 box plot + 個別 lap 点 overlay。Y = M:SS,CC。
    外れ値 = IQR robust rule（全表示 lap 横断・report-only flag）赤リング + 短ラベル。
    fastest lap は ★ + 注記。final-only / provisional-only / mixed の3モードで動作
    （df の列のみ使用・run_id 依存なし）。
    """
    plt, FuncFormatter = _mpl()
    series = []                          # (rec, [(x_jitter, lap_time, lap_no)])
    pos = 0
    for rec in _run_records(df):
        g = rec["_df"].sort_values("lap_no")
        vals = []
        for j, (_, row) in enumerate(g.iterrows()):
            t = _num(row.get("lap_time_s"))
            if t is None:
                continue
            vals.append((t, _num(row.get("lap_no")), j))
        if vals:
            pos += 1
            series.append((rec, pos, vals))
    if not series:
        return _text_page(tmpdir, "lap_time_dist", "Lap Time Distribution",
                          ["n/a — no valid lap times in selection"])
    fig, ax = plt.subplots(figsize=(11.5, 5.8))
    bp = ax.boxplot([[t for t, _ln, _j in vals] for _rec, _p, vals in series],
                    positions=[p for _rec, p, _v in series], widths=0.5,
                    showfliers=False, patch_artist=True,
                    medianprops=dict(color="#1A1A1A", lw=1.6))
    for i, box in enumerate(bp["boxes"]):
        c = RUN_PALETTE[i % len(RUN_PALETTE)]
        box.set_facecolor(c)
        box.set_alpha(0.15)
        box.set_edgecolor(c)
    all_t = [t for _rec, _p, vals in series for t, _ln, _j in vals]
    bounds = _iqr_bounds(all_t)
    fastest = min(all_t)
    fast_info = None
    outliers = []                        # (distance, x, t, label)
    for i, (rec, p, vals) in enumerate(series):
        c = RUN_PALETTE[i % len(RUN_PALETTE)]
        try:
            rtag = f"R{int(rec.get('run_no'))}"
        except (TypeError, ValueError):
            rtag = f"R{rec.get('run_no')}"
        for t, ln, j in vals:
            xj = p + ((j % 5) - 2) * 0.05          # 決定論 jitter（RNG 不使用）
            ax.plot([xj], [t], "o", ms=4.5, color=c, alpha=0.85, zorder=4)
            ln_txt = f"L{int(ln)}" if ln is not None else "L?"
            if t == fastest and fast_info is None:
                fast_info = (xj, t, f"{rtag} {ln_txt}")
            if bounds and (t < bounds[0] or t > bounds[1]):
                d = (bounds[0] - t) if t < bounds[0] else (t - bounds[1])
                outliers.append((d, xj, t, f"{rtag} {ln_txt} {format_lap_time(t)}"))
    if outliers:
        ax.scatter([o[1] for o in outliers], [o[2] for o in outliers], s=150,
                   facecolors="none", edgecolors="#D62728", linewidths=1.8, zorder=6)
        outliers.sort(key=lambda o: -o[0])
        for k, (_d, ox, ot, txt) in enumerate(outliers[:OUTLIER_LABEL_CAP]):
            dy = 9 if k % 2 == 0 else -14
            ax.annotate(txt, (ox, ot), xytext=(0, dy), textcoords="offset points",
                        ha="center", fontsize=7.5, color="#D62728",
                        fontweight="bold", zorder=7)
        if len(outliers) > OUTLIER_LABEL_CAP:
            ax.annotate(f"+{len(outliers) - OUTLIER_LABEL_CAP} more flagged",
                        (0.99, 0.02), xycoords="axes fraction", ha="right",
                        fontsize=8, color="#D62728")
    if fast_info:
        fx, ft, ftag = fast_info
        ax.plot([fx], [ft], marker="*", ms=16, color="#D4A017",
                markeredgecolor="black", markeredgewidth=0.6, zorder=7)
        ax.annotate(f"Fastest {format_lap_time(ft)} ({ftag})", (fx, ft),
                    xytext=(10, -14), textcoords="offset points", fontsize=9,
                    fontweight="bold", color="#8A6D00", zorder=7)
    lo_y, hi_y = ax.get_ylim()
    pad = (hi_y - lo_y) * 0.08 or 0.5
    ax.set_ylim(lo_y - pad, hi_y + pad)              # ラベル headroom
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: format_lap_time(v)))
    ax.set_ylabel("Lap time (M:SS,CC)")
    ax.set_xticks([p for _rec, p, _v in series])
    ax.set_xticklabels([run_short_label(rec) for rec, _p, _v in series],
                       rotation=30, ha="right", fontsize=9)
    ax.set_title("Lap Time Distribution  (box = Q1-Q3, ★ = fastest, red ring = outlier)")
    # PDF 経路にも規則開示が載るよう、注記を図内にも焼き込む（PPTX は slide note 併用）
    fig.text(0.5, 0.015,
             "Red ring = report-only lap-time outlier flag (IQR rule across all shown laps); "
             "laps are NOT removed - no DB/extraction change.",
             ha="center", fontsize=8, style="italic", color="#666666")
    fig.subplots_adjust(bottom=0.24)
    return _save(fig, tmpdir, "lap_time_dist")


def chart_lap_phase_progression(df, tmpdir, kind):
    """Lap by lap: phase position(kind='pos') or speed(kind='spd') 推移。

    3×2 small multiples（行=phase 色分け・列=F/R）。X=lap_no・Run毎の線。
    """
    plt, _ = _mpl()
    recs, extra = _capped_runs(df)
    fig, axes = plt.subplots(3, 2, figsize=(11.5, 6.6), sharex=True)
    for pr, phase in enumerate(PHASES):
        color = PHASE_COLORS[phase]
        if kind == "pos":
            cols = PHASE_POS[phase]
            side_tags = ("F pos", "R pos")
            unit = "mm"
        else:
            cols = (PHASE_SPD[phase]["F"][0], PHASE_SPD[phase]["R"][0])
            side_tags = (f"F {PHASE_SPD[phase]['F'][2]}", f"R {PHASE_SPD[phase]['R'][2]}")
            unit = "idx"
        for si, (col, tag) in enumerate(zip(cols, side_tags)):
            ax = axes[pr][si]
            has = False
            for i, rec in enumerate(recs):
                g = rec["_df"].sort_values("lap_no")
                if col not in g.columns:
                    continue
                xs, ys = [], []
                for _, row in g.iterrows():
                    ln, val = _num(row.get("lap_no")), _num(row.get(col))
                    if ln is not None and val is not None:
                        xs.append(int(ln))
                        ys.append(val)
                if xs:
                    has = True
                    ax.plot(xs, ys, "-o", ms=3, lw=1.3,
                            color=RUN_PALETTE[i % len(RUN_PALETTE)])
            ax.set_title(f"{phase} {tag} ({unit})", fontsize=10, color=color)
            if not has:
                ax.text(0.5, 0.5, "n/a", ha="center", va="center",
                        transform=ax.transAxes, color="#aaa", fontsize=12)
    for ax in axes[-1]:
        ax.set_xlabel("Lap No")
    # 共通凡例（Run 色）
    handles = [plt.Line2D([0], [0], color=RUN_PALETTE[i % len(RUN_PALETTE)], lw=2,
               label=run_short_label(rec)) for i, rec in enumerate(recs)]
    fig.legend(handles=handles, loc="lower center", ncol=min(len(handles), 6),
               frameon=False, fontsize=8, bbox_to_anchor=(0.5, -0.02))
    ttl = "Position" if kind == "pos" else "Suspension Speed"
    fig.suptitle(f"Lap-by-lap — Phase {ttl} Progression", fontsize=14, fontweight="bold")
    fig.subplots_adjust(hspace=0.42, wspace=0.22, bottom=0.12, top=0.90)
    return _save(fig, tmpdir, f"lap_phase_{kind}")


def chart_run_detail(rec, tmpdir):
    """1 Run の詳細: lap time / best差 と 3フェーズ position を small multiples。"""
    plt, FuncFormatter = _mpl()
    g = rec["_df"].sort_values("lap_no")
    s = lap_series(g)
    fig, axes = plt.subplots(2, 2, figsize=(11.0, 6.2))
    # (0,0) lap time
    ax = axes[0][0]
    xs = [ln for ln, t in zip(s["lap_no"], s["lap_time"]) if t is not None]
    ys = [t for t in s["lap_time"] if t is not None]
    if xs:
        ax.plot(xs, ys, "-o", color="#0078D4", ms=4)
        if s["best"] is not None:
            ax.plot([xs[ys.index(min(ys))]], [min(ys)], "*", color="#0078D4",
                    ms=14, markeredgecolor="black", markeredgewidth=0.5)
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: format_lap_time(v)))
    ax.set_title("Lap time (M:SS,CC)")
    # (0,1) best差
    ax = axes[0][1]
    dd = [(ln, d) for ln, d in zip(s["lap_no"], s["delta"]) if d is not None]
    if dd:
        ax.bar([a for a, _ in dd], [b for _, b in dd], color="#5D6D7E")
    ax.axhline(0, color="black", lw=0.8)
    ax.set_title("Δ to best lap (s)")
    # (1,0) position、(1,1) speed
    axp, axs = axes[1][0], axes[1][1]
    for phase in PHASES:
        col_f = PHASE_POS[phase][0]
        col_spd = PHASE_SPD[phase]["F"][0]
        c = PHASE_COLORS[phase]
        xs_p, ys_p, xs_s, ys_s = [], [], [], []
        for _, row in g.iterrows():
            ln = _num(row.get("lap_no"))
            if ln is None:
                continue
            vp, vs = _num(row.get(col_f)), _num(row.get(col_spd))
            if vp is not None:
                xs_p.append(int(ln)); ys_p.append(vp)
            if vs is not None:
                xs_s.append(int(ln)); ys_s.append(vs)
        if xs_p:
            axp.plot(xs_p, ys_p, "-o", ms=3, color=c, label=phase)
        if xs_s:
            axs.plot(xs_s, ys_s, "-o", ms=3, color=c, label=phase)
    axp.set_title("Front position by phase (mm)")
    axs.set_title("Front damping speed by phase (idx)")
    for ax in (axp, axs):
        ax.set_xlabel("Lap No")
        ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.22), ncol=3,
                  frameon=False, fontsize=8)
    fig.suptitle(f"Run Detail — {run_short_label(rec)}  ({rec['n_laps']} laps, best {format_lap_time(s['best'])})",
                 fontsize=13, fontweight="bold")
    fig.subplots_adjust(hspace=0.5, wspace=0.25, top=0.88, bottom=0.12)
    return _save(fig, tmpdir, f"run_detail_{rec['run_id']}")


# ── PPTX 組立 ─────────────────────────────────────────────────────────
def _pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN
    except Exception as exc:  # pragma: no cover
        raise ReportUnavailableError(f"python-pptx unavailable: {exc}") from exc


def _add_title(prs, Inches, Pt, RGBColor, PP_ALIGN, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(1.8))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    return slide


def _add_image_slide(prs, Inches, Pt, RGBColor, title, img_path, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)
    top = Inches(0.95)
    max_w, max_h = Inches(12.3), Inches(6.0)
    pic = slide.shapes.add_picture(str(img_path), Inches(0.5), top, height=max_h)
    if pic.width > max_w:                    # 幅超過時は縦横比維持で縮小
        ratio = max_w / pic.width
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int((prs.slide_width - pic.width) / 2)
    if note:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = note
        p2.font.size = Pt(9)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    return slide


def _style_cell(cell, Pt, RGBColor, text, *, bold=False, fill=None,
                size=10, align=None, color=(0x22, 0x22, 0x22)):
    from pptx.util import Pt as _Pt
    cell.text = str(text)
    para = cell.text_frame.paragraphs[0]
    para.font.size = _Pt(size)
    para.font.bold = bold
    para.font.color.rgb = RGBColor(*color)
    if align is not None:
        para.alignment = align
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill)


def _add_run_compare_table(prs, Inches, Pt, RGBColor, PP_ALIGN, df):
    """Run Compare Table（ヘッダ2行・単位・Braking/Apex/Exit 色分け・説明行）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Run Comparison Table"
    tp.font.size = Pt(22)
    tp.font.bold = True
    # 説明行
    nb = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.35))
    npp = nb.text_frame.paragraphs[0]
    npp.text = ("pos = sensor position [mm] | spd = " + SPEED_NOTE
                + " | n/a = not available")
    npp.font.size = Pt(9)
    npp.font.italic = True
    npp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    recs = _run_records(df)
    # 列: Run, Best, [Braking F pos, F dive], [Apex F pos, F dive], [Exit F reb, R |v|]
    col_specs = [
        ("Run", None), ("Best\n[M:SS,CC]", None),
        ("Brk F pos\n[mm]", "Braking"), ("Brk F-Dive\n[idx·deep]", "Braking"),
        ("Apex F pos\n[mm]", "Apex"), ("Apex F-Dive\n[idx·mid]", "Apex"),
        ("Exit F reb\n[idx]", "Exit"), ("Exit R |v|\n[idx]", "Exit"),
    ]
    ncols = len(col_specs)
    nrows = 1 + len(recs)
    gtab = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.2),
                                  Inches(12.3), Inches(0.4 + 0.32 * len(recs)))
    table = gtab.table
    for c, (name, phase) in enumerate(col_specs):
        fill = PHASE_FILL.get(phase) if phase else "D9D9D9"
        _style_cell(table.cell(0, c), Pt, RGBColor, name, bold=True, fill=fill,
                    size=9, align=PP_ALIGN.CENTER)
    for r, rec in enumerate(recs, start=1):
        g = rec["_df"]
        vals = [
            run_short_label(rec),
            format_lap_time(run_best(g)),
            _fmt_mm(_mean(g.get("brk_susf_avg"))),
            _fmt_idx(_mean(g.get("brk_f_dive_spd_avg"))),
            _fmt_mm(_mean(g.get("apex_susf_avg"))),
            _fmt_idx(_mean(g.get("apex_f_dive_spd_avg"))),
            _fmt_idx(_mean(g.get("ce_f_reb_spd_avg"))),
            _fmt_idx(_mean(g.get("ce_r_spd_avg"))),
        ]
        for c, (v, (_, phase)) in enumerate(zip(vals, col_specs)):
            fill = PHASE_FILL.get(phase) if phase else None
            _style_cell(table.cell(r, c), Pt, RGBColor, v, fill=fill, size=9,
                        align=(PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT))
    return slide


def _add_quality_table(prs, Inches, Pt, RGBColor, PP_ALIGN, df, filter_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Data Quality & Coverage"
    tp.font.size = Pt(22)
    tp.font.bold = True
    nb = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.5))
    npp = nb.text_frame.paragraphs[0]
    npp.text = ("Missing = laps with NULL value (not vehicle speed). "
                "0 is a real measurement, not missing. "
                "Exit metrics are structurally sparse (CORNER_EXIT).")
    npp.font.size = Pt(10)
    npp.font.italic = True
    npp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    rows = data_quality(df)
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.5), Inches(1.35),
                                   Inches(12.3), Inches(0.4 + 0.5 * len(rows))).table
    for c, name in enumerate(["Metric group", "Populated / Missing", "Meaning"]):
        _style_cell(table.cell(0, c), Pt, RGBColor, name, bold=True, fill="D9D9D9",
                    size=11, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        phase = next((p for p in PHASES if p.lower() in row["group"].lower()), None)
        fill = PHASE_FILL.get(phase) if phase else None
        miss = row["missing_pct"]
        pop_txt = (f"{row['populated']}/{row['total']} populated"
                   + (f" · Missing {miss:.0f}%" if miss is not None else " · n/a"))
        _style_cell(table.cell(r, 0), Pt, RGBColor, row["group"], bold=True, fill=fill, size=10)
        _style_cell(table.cell(r, 1), Pt, RGBColor, pop_txt, fill=fill, size=10)
        _style_cell(table.cell(r, 2), Pt, RGBColor, row["note"].split(" — ")[-1], fill=fill, size=9)
    if filter_note:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
        fb.text_frame.word_wrap = True
        fp = fb.text_frame.paragraphs[0]
        fp.text = filter_note
        fp.font.size = Pt(9)
        fp.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    return slide


def _fmt_mm(v):
    f = _num(v)
    return "n/a" if f is None else f"{f:.1f}"


def _fmt_idx(v):
    f = _num(v)
    return "n/a" if f is None else f"{f:.0f}"


# ── scope 正規化（PPTX/PDF 出力は英語のみ・Workbench の日本語 "全" を排除）──
ALL_SENTINELS = {"", "全", "全て", "全サーキット", "全ライダー", "全セッション",
                 "全ラウンド", "ALL", "ALL RIDERS", "ALL RUNS", "NA", None}


def _ascii_token(s):
    """ファイル名用 ASCII トークン（英数 + _ - ・大文字）。CJK 等は除去。空なら 'NA'。"""
    t = "".join(ch for ch in str(s)
                if ch.isascii() and (ch.isalnum() or ch in "_-"))
    return t.upper() or "NA"


def _resolve_one(raw, df, col, all_label):
    """(display, filename_token) を返す。'全' 等・空・複数は英語 all_label / 'ALL'。"""
    is_all = raw is None or str(raw).strip() in ALL_SENTINELS
    vals = (sorted(df[col].dropna().astype(str).unique().tolist())
            if col in df.columns else [])
    if is_all:
        if len(vals) == 1:
            return vals[0], _ascii_token(vals[0])
        disp = f"{all_label} ({', '.join(vals)})" if 1 < len(vals) <= 3 else all_label
        return disp, "ALL"
    return str(raw), _ascii_token(raw)


def _resolve_scope(scope, df):
    """Workbench から来る scope（日本語 '全' 含む）を英語表示 + ASCII トークンへ。"""
    scope = scope or {}
    c_d, c_t = _resolve_one(scope.get("circuit"), df, "circuit", "All circuits")
    r_d, r_t = _resolve_one(scope.get("rider"), df, "rider", "All riders")
    s_d, s_t = _resolve_one(scope.get("session"), df, "session", "All sessions")
    return {"circuit": c_d, "circuit_tok": c_t, "rider": r_d, "rider_tok": r_t,
            "session": s_d, "session_tok": s_t}


def chart_cover(resolved, summ, generated_human, tmpdir, provisional=False,
                mixed=False):
    """チーム提出用 Cover（Title / Subtitle / KPI cards / Scope / Phase legend）。

    16:9 全面。PPTX/PDF 両方でスライド1として使う。**英語のみ・CJK なし。**
    provisional=True でリボン + 注記4行を追加（§59 / 既存要素の座標は不変）。
    """
    plt, _ = _mpl()
    import matplotlib.patches as mpatches
    NAVY, INK, GREY = "#0B3D6B", "#1A1A1A", "#5B6672"
    CARD, EDGE, FOOT = "#F4F7FA", "#DCE3EA", "#8A94A0"
    fig = plt.figure(figsize=(13.333, 7.5))
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    ax.add_patch(mpatches.Rectangle((0, 0), 0.016, 1, color=NAVY, zorder=2))
    ax.add_patch(mpatches.Rectangle((0.05, 0.945), 0.17, 0.010, color=NAVY, zorder=2))
    ax.text(0.05, 0.87, "TS24 Suspension Performance Report",
            fontsize=31, fontweight="bold", color=INK, va="center")
    ax.text(0.05, 0.792,
            f"{resolved['circuit']}    ·    {resolved['rider']}    ·    {resolved['session']}",
            fontsize=17, color=GREY, va="center")
    ax.plot([0.05, 0.95], [0.752, 0.752], color=EDGE, lw=1.2)

    if provisional:
        # リボン（英語のみ / CJK=0）— 既存要素の座標は一切動かさない
        AMBER = "#B7791F"
        ax.add_patch(mpatches.FancyBboxPatch(
            (0.60, 0.930), 0.35, 0.058,
            boxstyle="round,pad=0,rounding_size=0.010",
            fc=AMBER, ec="none", zorder=4, mutation_aspect=1.777))
        ax.text(0.775, 0.959, "PROVISIONAL - SESSION DATA",
                fontsize=14.5, fontweight="bold", color="white",
                ha="center", va="center", zorder=5)
        if mixed:
            ax.text(0.95, 0.912, "Mixed final + provisional runs",
                    fontsize=9.5, fontweight="bold", color=AMBER,
                    ha="right", va="center", zorder=4)
        # メタデータ注記 4行（Scope/Phase カードと footer の間の空き帯・警告色）
        ny = 0.4975
        for line in ("Not final DB integration",
                     "Original setup data not merged",
                     "Run numbers are provisional",
                     "For race-weekend engineering review only"):
            ax.text(0.05, ny, line, fontsize=8.5, fontweight="bold",
                    color=AMBER, va="center", zorder=4)
            ny -= 0.0155

    kpis = [("Runs", str(summ["n_runs"])), ("Laps", str(summ["n_laps"])),
            ("Best lap", format_lap_time(summ["best"])),
            ("Median lap", format_lap_time(summ["median"]))]
    x0, x1, gap = 0.05, 0.95, 0.025
    w = (x1 - x0 - gap * (len(kpis) - 1)) / len(kpis)
    yb, h = 0.505, 0.20
    for i, (lab, val) in enumerate(kpis):
        x = x0 + i * (w + gap)
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, yb), w, h, boxstyle="round,pad=0,rounding_size=0.012",
            fc=CARD, ec=EDGE, lw=1.2, zorder=2, mutation_aspect=1.777))
        ax.text(x + w / 2, yb + h * 0.63, val, fontsize=26, fontweight="bold",
                color=NAVY, ha="center", va="center", zorder=3)
        ax.text(x + w / 2, yb + h * 0.24, lab, fontsize=13, color=GREY,
                ha="center", va="center", zorder=3)

    syb, sh = 0.085, 0.36
    ax.add_patch(mpatches.FancyBboxPatch((0.05, syb), 0.42, sh,
        boxstyle="round,pad=0,rounding_size=0.012", fc=CARD, ec=EDGE, lw=1.2,
        zorder=2, mutation_aspect=1.777))
    ax.text(0.075, syb + sh - 0.05, "Scope", fontsize=13, fontweight="bold",
            color=INK, va="center", zorder=3)
    ry = syb + sh - 0.12
    for lab, val in [("Circuit", resolved["circuit"]), ("Session", resolved["session"]),
                     ("Rider", resolved["rider"]), ("Generated", generated_human)]:
        ax.text(0.075, ry, lab, fontsize=12, color=GREY, va="center", zorder=3)
        ax.text(0.185, ry, str(val), fontsize=12.5, color=INK, va="center", zorder=3)
        ry -= 0.07

    ax.add_patch(mpatches.FancyBboxPatch((0.53, syb), 0.42, sh,
        boxstyle="round,pad=0,rounding_size=0.012", fc=CARD, ec=EDGE, lw=1.2,
        zorder=2, mutation_aspect=1.777))
    ax.text(0.555, syb + sh - 0.05, "Phase colours", fontsize=13, fontweight="bold",
            color=INK, va="center", zorder=3)
    ly = syb + sh - 0.13
    for name, desc in [("Braking", "entry / full braking"), ("Apex", "mid corner"),
                       ("Exit", "corner exit / drive")]:
        ax.add_patch(mpatches.FancyBboxPatch((0.555, ly - 0.020), 0.045, 0.040,
            boxstyle="round,pad=0,rounding_size=0.008", fc=PHASE_COLORS[name],
            ec="none", zorder=3, mutation_aspect=1.777))
        ax.text(0.625, ly, name, fontsize=13, fontweight="bold", color=INK,
                va="center", zorder=3)
        ax.text(0.755, ly, desc, fontsize=11, color=GREY, va="center", zorder=3)
        ly -= 0.083
    ax.text(0.05, 0.045, SPEED_NOTE, fontsize=9, style="italic", color=FOOT, va="center")
    return _save(fig, tmpdir, "cover", tight=False)


def _detect_provisional(df, provisional):
    """PROV_ run 自動検出（安全網・§59）。(provisional, mixed) を返す。

    呼び出し側が provisional=False でも PROV_ run を含めば強制昇格する
    （表記漏れ=提出事故の恒久防止）。final-only は明示 True を尊重。
    """
    prov_ids = {r for r in df["run_id"].astype(str) if r.startswith("PROV_")}
    if prov_ids:
        provisional = True
    mixed = bool(prov_ids) and len(prov_ids) < df["run_id"].nunique()
    return provisional, mixed


def _add_cover_slide(prs, Inches, img_path):
    """Cover 画像をスライド全面に配置（16:9 一致）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    return slide


def build_report_v2(df, run_ids=None, scope=None, out_dir=DEFAULT_OUT,
                    timestamp=None, provisional=False, lap_filter=True) -> Path:
    """フィルタ済み DataFrame から Report v2 PPTX を生成し保存パスを返す。"""
    import tempfile

    if run_ids:
        df = df[df["run_id"].isin(list(run_ids))].reset_index(drop=True)
    if df is None or df.empty:
        raise ReportUnavailableError("No lap_suspension rows for the selected filter/runs")
    provisional, mixed = _detect_provisional(df, provisional)
    df, _excluded = apply_lap_filter(df, lap_filter)
    if df is None or df.empty:
        raise ReportUnavailableError("No lap_suspension rows after lap filter")
    fnote = lap_filter_note(_excluded, lap_filter)

    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _pptx()
    plt, _ = _mpl()
    _style(plt)

    rs = _resolve_scope(scope, df)
    ts = timestamp or datetime.now().strftime("%Y%m%d_%H%M%S")
    gen_human = datetime.now().strftime("%Y-%m-%d %H:%M")
    summ = session_summary(df)
    recs = _run_records(df)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        # 1. Cover（チーム提出用・英語のみ）
        _add_cover_slide(prs, Inches, chart_cover(rs, summ, gen_human, tmp,
                                                  provisional=provisional,
                                                  mixed=mixed))
        # 2. Data quality（page 2 に lap-filter 開示を含む）
        _add_quality_table(prs, Inches, Pt, RGBColor, PP_ALIGN, df, filter_note=fnote)
        # 3. Run overview
        _add_image_slide(prs, Inches, Pt, RGBColor, "Run Overview",
                         chart_run_overview(df, tmp),
                         note="Lap time formatted M:SS,CC — labels placed outside bars")
        # 4-6. Phase summaries
        for phase in PHASES:
            _add_image_slide(prs, Inches, Pt, RGBColor, f"{phase} Phase Summary",
                             chart_phase_summary(df, phase, tmp), note=SPEED_WINDOW_NOTE)
        # 6b. All laps phase trend & outliers（phase summary の後・lap-by-lap の前）
        _add_image_slide(prs, Inches, Pt, RGBColor, "All Laps Phase Trend & Outliers",
                         chart_all_laps_phase_trend(df, tmp), note=TREND_OUTLIER_NOTE)
        # 7. Lap-by-lap time progression
        _add_image_slide(prs, Inches, Pt, RGBColor, "Lap-by-lap: Lap Time",
                         chart_lap_time_progression(df, tmp),
                         note="★ = run best lap. Legend outside plot to avoid overlap.")
        # 7b. Lap time distribution（lap-time progression の直後）
        _add_image_slide(prs, Inches, Pt, RGBColor, "Lap Time Distribution",
                         chart_lap_time_distribution(df, tmp), note=DIST_NOTE)
        # 8. Lap-by-lap position progression
        _add_image_slide(prs, Inches, Pt, RGBColor, "Lap-by-lap: Phase Position",
                         chart_lap_phase_progression(df, tmp, "pos"),
                         note="Braking=red / Apex=blue / Exit=green. n/a = not available.")
        # 9. Lap-by-lap speed progression
        _add_image_slide(prs, Inches, Pt, RGBColor, "Lap-by-lap: Suspension Speed",
                         chart_lap_phase_progression(df, tmp, "spd"), note=SPEED_WINDOW_NOTE)
        # 10..N Run detail
        for rec in recs[:RUN_DETAIL_CAP]:
            _add_image_slide(prs, Inches, Pt, RGBColor,
                             f"Run Detail — {run_short_label(rec)}",
                             chart_run_detail(rec, tmp),
                             note=SPEED_NOTE)
        if len(recs) > RUN_DETAIL_CAP:
            _add_title(prs, Inches, Pt, RGBColor, PP_ALIGN, "Run detail pages capped",
                       f"{len(recs) - RUN_DETAIL_CAP} more run(s) not shown as detail pages "
                       f"(cap {RUN_DETAIL_CAP}). See Run Comparison Table for all runs.")
        # Run compare table
        _add_run_compare_table(prs, Inches, Pt, RGBColor, PP_ALIGN, df)
        # Data limits
        _add_title(prs, Inches, Pt, RGBColor, PP_ALIGN, "Data limits & notes",
                   SPEED_WINDOW_NOTE + "  " + PEAK_NOTE + "  " + OHLINS_NOTE
                   + "  Exit (CORNER_EXIT) is structurally sparse. 0 != missing. "
                   "See report audit reports/report_v2_feedback_audit_20260708.md (§66).")

        prov_tok = "PROVISIONAL_" if provisional else ""
        out = out_dir / (f"suspension_report_v2_{rs['circuit_tok']}_{rs['rider_tok']}_"
                         f"{rs['session_tok']}_{prov_tok}{ts}.pptx")
        prs.save(str(out))
    return out


def _first(df, col):
    if col in df.columns:
        vals = df[col].dropna().unique().tolist()
        if len(vals) == 1:
            return vals[0]
        return "ALL" if len(vals) > 1 else None
    return None


# ── PDF プレビュー（単一ファイル・macOS Preview で開ける）─────────────
def _text_page(tmpdir, name, title, lines, color="#1F1F1F"):
    """タイトル/注記のテキストページを PNG で描く。"""
    plt, _ = _mpl()
    fig = plt.figure(figsize=(11.5, 6.5))
    fig.text(0.06, 0.86, title, fontsize=24, fontweight="bold", color=color)
    y = 0.72
    for ln in lines:
        fig.text(0.06, y, ln, fontsize=13, color="#333333")
        y -= 0.09
    fig.patch.set_facecolor("white")
    return _save(fig, tmpdir, name)


def _table_page(tmpdir, name, title, col_labels, rows, row_fills, note=None,
                header_fills=None, col_widths=None):
    """matplotlib テーブルページ（行=row_fills / ヘッダ列=header_fills で塗る）。"""
    plt, _ = _mpl()
    fig, ax = plt.subplots(figsize=(11.5, 6.5))
    ax.axis("off")
    ax.set_title(title, fontsize=18, fontweight="bold", loc="left", pad=18)
    tbl = ax.table(cellText=rows, colLabels=col_labels, loc="center",
                   cellLoc="center", colWidths=col_widths)
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.7)
    ncol = len(col_labels)
    for c in range(ncol):                       # ヘッダ
        cell = tbl[0, c]
        hf = header_fills[c] if header_fills and header_fills[c] else "D9D9D9"
        cell.set_facecolor("#" + hf)
        cell.set_text_props(weight="bold")
    for r, fill in enumerate(row_fills, start=1):
        for c in range(ncol):
            cell = tbl[r, c]
            if fill:
                cell.set_facecolor("#" + fill)
            cell.set_text_props(ha="left" if c == ncol - 1 else "center")
    if note:
        fig.text(0.06, 0.06, note, fontsize=9, style="italic", color="#888888")
    fig.patch.set_facecolor("white")
    return _save(fig, tmpdir, name)


def _quality_page(df, tmpdir, filter_note=None):
    rows, fills = [], []
    for row in data_quality(df):
        phase = next((p for p in PHASES if p.lower() in row["group"].lower()), None)
        fills.append(PHASE_FILL.get(phase))
        miss = row["missing_pct"]
        pop = (f"{row['populated']}/{row['total']} populated"
               + (f" · Missing {miss:.0f}%" if miss is not None else " · n/a"))
        rows.append([row["group"], pop, row["note"].split(" — ")[-1]])
    note = "0 is a real measurement, not missing. Exit metrics are structurally sparse (CORNER_EXIT)."
    if filter_note:
        note = filter_note + "\n" + note
    return _table_page(tmpdir, "t_quality", "Data Quality & Coverage",
                       ["Metric group", "Populated / Missing", "Meaning"], rows, fills, note=note)


def _compare_page(df, tmpdir):
    specs = [("Run", None), ("Best\n[M:SS,CC]", None),
             ("Brk F pos\n[mm]", "Braking"), ("Brk F-Dive\n[idx·deep]", "Braking"),
             ("Apex F pos\n[mm]", "Apex"), ("Apex F-Dive\n[idx·mid]", "Apex"),
             ("Exit F reb\n[idx]", "Exit"), ("Exit R |v|\n[idx]", "Exit")]
    labels = [s[0] for s in specs]
    header_fills = [PHASE_FILL.get(ph) for _, ph in specs]
    recs = _run_records(df)
    multi_sess = df["session"].nunique() > 1 if "session" in df.columns else False
    rows = []
    for rec in recs:
        g = rec["_df"]
        try:
            tag = f"R{int(rec.get('run_no'))}"
        except (TypeError, ValueError):
            tag = f"R{rec.get('run_no')}"
        if multi_sess:
            tag = f"{rec.get('session')} {tag}"
        rows.append([tag, format_lap_time(run_best(g)),
                     _fmt_mm(_mean(g.get("brk_susf_avg"))), _fmt_idx(_mean(g.get("brk_f_dive_spd_avg"))),
                     _fmt_mm(_mean(g.get("apex_susf_avg"))), _fmt_idx(_mean(g.get("apex_f_dive_spd_avg"))),
                     _fmt_idx(_mean(g.get("ce_f_reb_spd_avg"))), _fmt_idx(_mean(g.get("ce_r_spd_avg")))])
    widths = [0.14, 0.13] + [0.1225] * 6      # Run/Best を広め・phase 6列を均等
    return _table_page(tmpdir, "t_compare", "Run Comparison Table", labels, rows,
                       [None] * len(rows), note="pos = position [mm] | " + SPEED_NOTE,
                       header_fills=header_fills, col_widths=widths)


def _png_uniform(path, W=1600, H=925):
    """PNG を白キャンバス(W×H)へ縦横比維持で中央配置し RGB 化。"""
    from PIL import Image
    im = Image.open(path).convert("RGBA")
    bg = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    r = min(W / im.width, H / im.height)
    im = im.resize((max(1, int(im.width * r)), max(1, int(im.height * r))))
    bg.paste(im, ((W - im.width) // 2, (H - im.height) // 2), im)
    return bg.convert("RGB")


def build_report_pdf(df, run_ids=None, scope=None, out_dir=DEFAULT_OUT,
                     timestamp=None, provisional=False, lap_filter=True) -> Path:
    """PPTX と同じ内容を単一 PDF で出力（macOS Preview で開ける）。"""
    import tempfile
    from PIL import Image

    if run_ids:
        df = df[df["run_id"].isin(list(run_ids))].reset_index(drop=True)
    if df is None or df.empty:
        raise ReportUnavailableError("No lap_suspension rows for the selected filter/runs")
    provisional, mixed = _detect_provisional(df, provisional)
    df, _excluded = apply_lap_filter(df, lap_filter)
    if df is None or df.empty:
        raise ReportUnavailableError("No lap_suspension rows after lap filter")
    fnote = lap_filter_note(_excluded, lap_filter)
    plt, _ = _mpl()
    _style(plt)
    rs = _resolve_scope(scope, df)
    ts = timestamp or datetime.now().strftime("%Y%m%d_%H%M%S")
    gen_human = datetime.now().strftime("%Y-%m-%d %H:%M")
    summ = session_summary(df)
    recs = _run_records(df)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        pages = []
        pages.append(chart_cover(rs, summ, gen_human, tmp,   # 1. Cover（英語のみ）
                                 provisional=provisional, mixed=mixed))
        pages.append(_quality_page(df, tmp, filter_note=fnote))
        pages.append(chart_run_overview(df, tmp))
        for ph in PHASES:
            pages.append(chart_phase_summary(df, ph, tmp))
        pages.append(chart_all_laps_phase_trend(df, tmp))     # 指示書§2（PPTX と同順）
        pages.append(chart_lap_time_progression(df, tmp))
        pages.append(chart_lap_time_distribution(df, tmp))    # 指示書§3（PPTX と同順）
        pages.append(chart_lap_phase_progression(df, tmp, "pos"))
        pages.append(chart_lap_phase_progression(df, tmp, "spd"))
        for rec in recs[:RUN_DETAIL_CAP]:
            pages.append(chart_run_detail(rec, tmp))
        pages.append(_compare_page(df, tmp))
        pages.append(_text_page(tmp, "p_limits", "Data limits & notes", [
            "Susp speed = relative damping-speed index (mm/s, uncalibrated), NOT vehicle speed.",
            "It is the MEAN velocity within each phase window, not the peak brake dive-in rate.",
            "Braking window = deep/settled fork (90-130mm); Apex = mid-stroke (50-100mm).",
            "Do not read Apex>Braking as 'the front dives faster at apex' (see audit §66, 2026-07-08).",
            "avg needs n>=5 laps. Peak columns not shown (brk=legacy MAX vs others p95, not comparable).",
            "NOT comparable to Ohlins low/high-speed C/R (force-vs-shaft-velocity) — different quantity.",
            "Exit (CORNER_EXIT) is structurally sparse. 0 != missing."]))
        imgs = [_png_uniform(p) for p in pages]
        prov_tok = "PROVISIONAL_" if provisional else ""
        out = out_dir / (f"suspension_report_v2_{rs['circuit_tok']}_{rs['rider_tok']}_"
                         f"{rs['session_tok']}_{prov_tok}{ts}.pdf")
        imgs[0].save(str(out), save_all=True, append_images=imgs[1:], resolution=150)
    return out


# ── CLI ────────────────────────────────────────────────────────────────
def main(argv=None):
    import argparse
    ap = argparse.ArgumentParser(description="TS24 Suspension Report v2 generator")
    ap.add_argument("--db", default=str(DEFAULT_DB))
    ap.add_argument("--circuit")
    ap.add_argument("--rider")
    ap.add_argument("--session")
    ap.add_argument("--runs", nargs="*", help="run_id list (default: all in filter)")
    ap.add_argument("--out", default=str(DEFAULT_OUT))
    ap.add_argument("--timestamp")
    ap.add_argument("--pdf", action="store_true",
                    help="単一 PDF も出力（macOS Preview で開ける）")
    ap.add_argument("--no-lap-filter", action="store_true",
                    help="slow/out-lap の report-only フィルタを無効化（既定=有効・§66）")
    args = ap.parse_args(argv)

    df = load_lap_suspension(args.db, args.circuit, args.rider, args.session, args.runs)
    if df.empty:
        print("No rows for filter.")
        return 1
    scope = {"circuit": args.circuit, "rider": args.rider, "session": args.session}
    lap_filter = not args.no_lap_filter
    out = build_report_v2(df, run_ids=args.runs, scope=scope, out_dir=args.out,
                          timestamp=args.timestamp, lap_filter=lap_filter)
    print(f"Report written: {out}")
    if args.pdf:
        pdf = build_report_pdf(df, run_ids=args.runs, scope=scope, out_dir=args.out,
                               timestamp=args.timestamp, lap_filter=lap_filter)
        print(f"PDF preview:    {pdf}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
